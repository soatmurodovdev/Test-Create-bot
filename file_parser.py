# Murodjon Soatmurodov tomonidan yaratilgan
"""
Turli fayl formatlaridan matnni ajratib olish uchun modul.
Qo'llab-quvvatlanadi: PDF, DOCX, XLSX/XLS, TXT, CSV, RTF, PPTX, EPUB
"""
import csv


def extract_text(file_path: str, ext: str) -> str:
    ext = ext.lower().lstrip(".")
    if ext == "pdf":
        return _extract_pdf(file_path)
    if ext == "docx":
        return _extract_docx(file_path)
    if ext == "doc":
        raise ValueError(
            "Eski .doc format qo'llab-quvvatlanmaydi. Iltimos, faylni .docx formatga saqlab qayta yuboring."
        )
    if ext in ("xlsx", "xls"):
        return _extract_xlsx(file_path)
    if ext == "txt":
        return _read_txt(file_path)
    if ext == "csv":
        return _extract_csv(file_path)
    if ext == "rtf":
        return _extract_rtf(file_path)
    if ext == "pptx":
        return _extract_pptx(file_path)
    if ext == "epub":
        return _extract_epub(file_path)
    raise ValueError(f"Qo'llab-quvvatlanmaydigan format: .{ext}")


def _read_txt(path):
    for enc in ("utf-8", "cp1251", "latin-1"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="ignore")


def _extract_pdf(path):
    import pdfplumber

    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append(f"\n--- Sahifa {i} ---\n{text}")
    return "".join(pages)


def _extract_docx(path):
    import docx

    d = docx.Document(path)
    parts = []
    for p in d.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(path):
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"\n--- Varaq: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            values = [str(c) for c in row if c is not None]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _extract_csv(path):
    parts = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            parts.append(" | ".join(row))
    return "\n".join(parts)


def _extract_rtf(path):
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        raise ValueError("RTF o'qish uchun 'striprtf' kutubxonasi kerak: pip install striprtf")
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return rtf_to_text(f.read())


def _extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n--- Slayd {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_epub(path):
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ValueError("EPUB o'qish uchun 'ebooklib' va 'beautifulsoup4' kerak.")
    book = epub.read_epub(path)
    parts = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n").strip()
            if text:
                parts.append(text)
    return "\n".join(parts)
